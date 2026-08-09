"""Generate tiny, original test documents; no third-party sample assets are copied."""

from pathlib import Path


OUTPUT = Path(__file__).parent / "generated"


def main() -> None:
    try:
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit("Install the project Docling extra before generating samples.") from exc

    OUTPUT.mkdir(exist_ok=True)

    _write_pdf(OUTPUT / "invoice.pdf")

    docx = Document()
    docx.add_heading("Invoice SAMPLE-001", level=1)
    docx.add_paragraph("Total: 100.00 USD")
    docx.save(OUTPUT / "invoice.docx")

    pptx = Presentation()
    slide = pptx.slides.add_slide(pptx.slide_layouts[1])
    slide.shapes.title.text = "Invoice SAMPLE-001"
    slide.placeholders[1].text = "Total: 100.00 USD"
    pptx.save(OUTPUT / "invoice.pptx")

    xlsx = Workbook()
    sheet = xlsx.active
    sheet.title = "Invoice"
    sheet.append(["Invoice", "SAMPLE-001"])
    sheet.append(["Total", 100.00])
    xlsx.save(OUTPUT / "invoice.xlsx")

    print(f"Created samples in {OUTPUT}")


def _write_pdf(path: Path) -> None:
    """Write a valid one-page PDF using only the standard library."""

    stream = b"BT /F1 16 Tf 72 760 Td (Invoice SAMPLE-001) Tj 0 -22 Td (Total: 100.00 USD) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    content = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, obj in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{number} 0 obj\n".encode() + obj + b"\nendobj\n")
    xref = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    content.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode())
    content.extend(
        f"trailer << /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    )
    path.write_bytes(content)


if __name__ == "__main__":
    main()
